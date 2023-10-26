import pptx

# Ouvrir le fichier pptx
presentation = pptx.Presentation("my_presentation.pptx")

# Parcourir les diapositives
for slide in presentation.slides:

    # Afficher le titre de la diapositive
    print(slide.title)

    # Afficher les images de la diapositive
    for image in slide.shapes:
        if image.shape_type == pptx.enum.shapes.ShapeType.IMAGE:
            print(image.image.filename)

    # Afficher le texte de la diapositive
    for shape in slide.shapes:
        if shape.shape_type == pptx.enum.shapes.ShapeType.TEXTBOX:
            print(shape.text)